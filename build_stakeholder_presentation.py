#!/usr/bin/env python3
"""Run full analysis on the Excel dataset and build a 5-slide stakeholder presentation."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import matplotlib.pyplot as plt
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt


DATA_FILE = Path("IMS_Test_Data 1.xlsx")
OUTPUT_DIR = Path("presentation_output")
CHARTS_DIR = OUTPUT_DIR / "charts"
PPT_FILE = OUTPUT_DIR / "MMC_Dataset_Stakeholder_Presentation.pptx"
SUMMARY_FILE = OUTPUT_DIR / "analysis_summary.md"


@dataclass
class DatasetStats:
    n_rows: int
    n_cols: int
    median_age: float
    age_iqr_low: float
    age_iqr_high: float
    female_pct: float
    male_pct: float


def _value_counts(df: pd.DataFrame, col: str, top_n: int = 10) -> pd.DataFrame:
    return (
        df[col]
        .astype("string")
        .fillna("<NA>")
        .value_counts(dropna=False)
        .head(top_n)
        .rename_axis(col)
        .rename("count")
        .to_frame()
    )


def _multi_select_rate(df: pd.DataFrame, col: str) -> int:
    # Multi-select columns are populated only when selected.
    return int(df[col].notna().sum())


def _save_barh(series: pd.Series, title: str, xlabel: str, output_path: Path, color: str) -> None:
    fig, ax = plt.subplots(figsize=(10, 5.5))
    series.sort_values(ascending=True).plot(kind="barh", ax=ax, color=color)
    ax.set_title(title, fontsize=14, pad=12)
    ax.set_xlabel(xlabel)
    ax.set_ylabel("")
    for i, v in enumerate(series.sort_values(ascending=True).values):
        ax.text(v + max(series.max() * 0.01, 0.8), i, str(int(v)), va="center", fontsize=9)
    fig.tight_layout()
    fig.savefig(output_path, dpi=200)
    plt.close(fig)


def _save_age_histogram(age_series: pd.Series, output_path: Path) -> None:
    fig, ax = plt.subplots(figsize=(10, 5.5))
    sns.histplot(age_series.dropna(), bins=14, kde=True, color="#006D77", ax=ax)
    ax.set_title("Age Distribution", fontsize=14, pad=12)
    ax.set_xlabel("Age")
    ax.set_ylabel("Respondents")
    fig.tight_layout()
    fig.savefig(output_path, dpi=200)
    plt.close(fig)


def _add_slide_title(slide, text: str) -> None:
    title = slide.shapes.title
    if title is None:
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6)).text_frame
        title.text = text
    else:
        title.text = text
    if hasattr(title, "text_frame"):
        p = title.text_frame.paragraphs[0]
        p.font.size = Pt(30)
        p.font.bold = True


def _add_bullets(slide, bullets: Iterable[str], left: float, top: float, width: float, height: float, size: int = 20) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    for idx, line in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(size)


def _add_picture(slide, image_path: Path, left: float, top: float, width: float) -> None:
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))


def build() -> None:
    if not DATA_FILE.exists():
        raise FileNotFoundError(f"Missing dataset: {DATA_FILE}")

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    CHARTS_DIR.mkdir(parents=True, exist_ok=True)

    sns.set_theme(style="whitegrid")

    df = pd.read_excel(DATA_FILE, sheet_name="Sheet1")

    age_col = "How old are you?"
    sex_col = "[Enumerator observation: Sex]"
    nationality_col = "What is your country of nationality?"
    legal_col = "What is your current migration/legal status?"
    location_col = "What country is the respondent in right now?"
    education_col = "What is the highest level of education you have completed?"
    income_col = "Were you making money in the 12 months before you left your country of departure?"
    duration_col = "_duration (minutes)"

    incident_cols = [
        "Q118/Robbery",
        "Q118/Physical violence",
        "Q118/Witnessed death",
        "Q118/Injury / ill-health from harsh conditions",
        "Q118/Detention",
        "Q118/Bribery/extortion",
        "Q118/Trafficking and exploitation",
    ]

    protection_cols = [
        "Q119/Travel in a group",
        "Q119/Avoid large cities",
        "Q119/Look for information and follow recommendations",
        "Q119/Plan my journey carefully",
        "Q119/Keep in regular contact with family",
        "Q119/Use safer method of transport",
        "Q119/Stop in places with trusted contacts",
    ]

    transport_cols = [
        "Q73/Walk",
        "Q73/Bus",
        "Q73/Truck",
        "Q73/Boat",
        "Q73/Car /pick-up",
        "Q73/Motorbike",
        "Q73/Aeroplane",
    ]

    n = len(df)

    age_series = pd.to_numeric(df[age_col], errors="coerce")
    sex_counts = _value_counts(df, sex_col)
    nat_counts = _value_counts(df, nationality_col, top_n=6)
    legal_counts = _value_counts(df, legal_col, top_n=6)
    location_counts = _value_counts(df, location_col, top_n=5)
    education_counts = _value_counts(df, education_col, top_n=6)
    income_counts = _value_counts(df, income_col, top_n=5)

    incident_counts = pd.Series({col.replace("Q118/", ""): _multi_select_rate(df, col) for col in incident_cols})
    protection_counts = pd.Series({col.replace("Q119/", ""): _multi_select_rate(df, col) for col in protection_cols})
    transport_counts = pd.Series({col.replace("Q73/", ""): _multi_select_rate(df, col) for col in transport_cols})

    female_pct = (sex_counts.loc["Female", "count"] / n * 100) if "Female" in sex_counts.index else 0.0
    male_pct = (sex_counts.loc["Male", "count"] / n * 100) if "Male" in sex_counts.index else 0.0

    stats = DatasetStats(
        n_rows=n,
        n_cols=df.shape[1],
        median_age=float(age_series.median()),
        age_iqr_low=float(age_series.quantile(0.25)),
        age_iqr_high=float(age_series.quantile(0.75)),
        female_pct=float(female_pct),
        male_pct=float(male_pct),
    )

    age_chart = CHARTS_DIR / "age_hist.png"
    nat_chart = CHARTS_DIR / "nationality_top.png"
    incident_chart = CHARTS_DIR / "incident_counts.png"
    protection_chart = CHARTS_DIR / "protection_counts.png"
    transport_chart = CHARTS_DIR / "transport_counts.png"

    _save_age_histogram(age_series, age_chart)
    _save_barh(nat_counts["count"], "Top Nationalities", "Respondents", nat_chart, "#005F73")
    _save_barh(incident_counts, "Reported Incidents on Journey", "Respondents", incident_chart, "#AE2012")
    _save_barh(protection_counts, "Top Protection Behaviors", "Respondents", protection_chart, "#0A9396")
    _save_barh(transport_counts, "Transport Modes Used", "Respondents", transport_chart, "#3A86FF")

    prs = Presentation()

    # Slide 1: Title and study overview.
    s1 = prs.slides.add_slide(prs.slide_layouts[5])
    _add_slide_title(s1, "Migration Journey Dataset: Key Findings")
    _add_bullets(
        s1,
        [
            "Objective: summarize respondent profile and journey risks for stakeholders.",
            f"Sample size: {stats.n_rows} records, {stats.n_cols} variables.",
            "Tooling: Python (pandas, matplotlib, seaborn, python-pptx).",
            "Source file: IMS_Test_Data 1.xlsx (Sheet1).",
        ],
        left=0.8,
        top=1.5,
        width=11.4,
        height=3.8,
        size=24,
    )

    # Slide 2: Demographics.
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    _add_slide_title(s2, "1) Respondent Demographics")
    _add_bullets(
        s2,
        [
            f"Median age: {stats.median_age:.1f} years (IQR {stats.age_iqr_low:.0f}-{stats.age_iqr_high:.0f}).",
            f"Sex split: Female {stats.female_pct:.1f}%, Male {stats.male_pct:.1f}%.",
            "Education mostly primary or secondary; low tertiary completion.",
        ],
        left=0.6,
        top=1.0,
        width=6.2,
        height=2.0,
        size=18,
    )
    _add_picture(s2, age_chart, left=0.6, top=3.0, width=5.8)
    _add_picture(s2, nat_chart, left=6.6, top=1.4, width=5.9)

    # Slide 3: Migration context.
    s3 = prs.slides.add_slide(prs.slide_layouts[5])
    _add_slide_title(s3, "2) Migration & Legal Context")
    top_legal = ", ".join([f"{idx} ({int(row['count'])})" for idx, row in legal_counts.head(3).iterrows()])
    top_location = ", ".join([f"{idx} ({int(row['count'])})" for idx, row in location_counts.head(2).iterrows()])
    top_edu = ", ".join([f"{idx} ({int(row['count'])})" for idx, row in education_counts.head(3).iterrows()])
    _add_bullets(
        s3,
        [
            f"Current location concentration: {top_location}.",
            f"Top legal statuses: {top_legal}.",
            f"Top education levels: {top_edu}.",
            f"Income before departure: No {int(income_counts.loc['No','count'])}, Yes {int(income_counts.loc['Yes','count'])}.",
        ],
        left=0.8,
        top=1.4,
        width=11.0,
        height=4.6,
        size=20,
    )

    # Slide 4: Risks and journey methods.
    s4 = prs.slides.add_slide(prs.slide_layouts[5])
    _add_slide_title(s4, "3) Journey Methods and Exposure to Harm")
    _add_bullets(
        s4,
        [
            "Walking and bus travel dominate mobility patterns.",
            "Robbery and physical violence are the most frequently reported incidents.",
            "Multiple respondents also report witnessed death and health impacts.",
        ],
        left=0.6,
        top=1.0,
        width=11.5,
        height=1.8,
        size=18,
    )
    _add_picture(s4, transport_chart, left=0.6, top=2.2, width=5.8)
    _add_picture(s4, incident_chart, left=6.6, top=2.2, width=5.8)

    # Slide 5: Protection and actions.
    s5 = prs.slides.add_slide(prs.slide_layouts[5])
    _add_slide_title(s5, "4) Protection Behaviors and Priority Actions")
    _add_picture(s5, protection_chart, left=0.6, top=1.5, width=5.8)
    _add_bullets(
        s5,
        [
            "Most common strategy is traveling in groups.",
            "A large subgroup reports taking no protective action.",
            "Priority: scale trusted-information channels and route safety guidance.",
            "Priority: targeted support where robbery and violence rates are highest.",
            "Priority: strengthen case management for high-risk groups.",
        ],
        left=6.6,
        top=1.5,
        width=5.6,
        height=4.8,
        size=18,
    )

    prs.save(PPT_FILE)

    summary_lines = [
        "# Statistical Analysis Summary",
        "",
        "Software / language used:",
        "- Python 3.x",
        "- pandas, seaborn, matplotlib, python-pptx",
        "",
        "Dataset:",
        f"- File: {DATA_FILE.name}",
        f"- Rows: {stats.n_rows}",
        f"- Columns: {stats.n_cols}",
        "",
        "Key demographics:",
        f"- Median age: {stats.median_age:.1f} years",
        f"- Female: {stats.female_pct:.1f}%",
        f"- Male: {stats.male_pct:.1f}%",
        f"- Top nationality: {nat_counts.index[0]} ({int(nat_counts.iloc[0]['count'])})",
        "",
        "Key findings:",
        f"- Most common legal status: {legal_counts.index[0]} ({int(legal_counts.iloc[0]['count'])})",
        f"- Top incident: {incident_counts.sort_values(ascending=False).index[0]} ({int(incident_counts.max())})",
        f"- Top protection behavior: {protection_counts.sort_values(ascending=False).index[0]} ({int(protection_counts.max())})",
        f"- Most used transport mode: {transport_counts.sort_values(ascending=False).index[0]} ({int(transport_counts.max())})",
        "",
        "Outputs generated:",
        f"- Presentation: {PPT_FILE}",
        f"- Charts directory: {CHARTS_DIR}",
    ]
    SUMMARY_FILE.write_text("\n".join(summary_lines), encoding="utf-8")

    print(f"Created presentation: {PPT_FILE}")
    print(f"Created summary: {SUMMARY_FILE}")


if __name__ == "__main__":
    build()
